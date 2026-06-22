"""
PowerPoint export for GA4/GTM audit results.
"""

from datetime import datetime
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from core import CheckResult, score_checks

# Maximum check rows per slide. Audits with more checks paginate onto
# additional "(cont.)" slides so nothing is ever silently dropped.
_ROWS_PER_SLIDE = 9

# Cap very long detail strings (e.g. the network session timeline) so a
# single cell can't blow up the table height. Full text stays in the terminal.
_MAX_DETAIL_CHARS = 260

_HEADER_FILL = RGBColor(33, 33, 33)
_HEADER_TEXT = RGBColor(255, 255, 255)
_DETAIL_TEXT = RGBColor(80, 80, 80)
_ROW_FILL_A = RGBColor(255, 255, 255)
_ROW_FILL_B = RGBColor(242, 242, 242)

# Colors
_COLOR_PASS = RGBColor(34, 177, 76)      # green
_COLOR_FAIL = RGBColor(237, 28, 36)      # red
_COLOR_SKIP = RGBColor(255, 192, 0)      # orange/yellow
_COLOR_INFO = RGBColor(0, 150, 200)      # cyan
_COLOR_CRITICAL = RGBColor(237, 28, 36)  # red (bold)
_COLOR_HIGH = RGBColor(237, 28, 36)      # red
_COLOR_MEDIUM = RGBColor(255, 192, 0)    # yellow
_COLOR_LOW = RGBColor(192, 192, 192)     # gray

_SEVERITY_COLOR = {
    "CRITICAL": _COLOR_CRITICAL,
    "HIGH": _COLOR_HIGH,
    "MEDIUM": _COLOR_MEDIUM,
    "LOW": _COLOR_LOW,
    "INFO": _COLOR_INFO,
}

_STATUS_COLOR = {
    "PASS": _COLOR_PASS,
    "FAIL": _COLOR_FAIL,
    "SKIP": _COLOR_SKIP,
}


def _score_color_rgb(pct: float) -> RGBColor:
    if pct >= 80:
        return _COLOR_PASS
    if pct >= 50:
        return _COLOR_MEDIUM
    return _COLOR_FAIL


def _add_title_slide(prs: Presentation, site_url: str) -> None:
    """Add opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(33, 33, 33)  # dark gray

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "GA4 / GTM Audit Results"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7), Inches(9), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = site_url
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(9), Inches(0.5)
    )
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)


def _status_of(check: CheckResult) -> Tuple[str, RGBColor]:
    if check.skipped:
        return "SKIP", _STATUS_COLOR["SKIP"]
    if check.passed:
        return "PASS", _STATUS_COLOR["PASS"]
    return "FAIL", _STATUS_COLOR["FAIL"]


def _set_cell(
    cell,
    text: str,
    *,
    size: int = 9,
    color: RGBColor = RGBColor(0, 0, 0),
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    fill: RGBColor = None,
) -> None:
    """Write and style a single table cell."""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)

    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    # Style via the paragraph font (covers the implicit run created by .text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_audit_slide(
    prs: Presentation,
    title: str,
    checks: List[CheckResult],
    stats: Tuple[int, int, int, int],
    score: Tuple[int, int, float],
) -> None:
    """Render one slide: title, score, stats line, and a table of the given checks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6.5), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(33, 33, 33)

    # Score badge (top-right)
    s_passed, s_total, s_pct = score
    score_box = slide.shapes.add_textbox(Inches(7.2), Inches(0.3), Inches(2.4), Inches(0.6))
    p = score_box.text_frame.paragraphs[0]
    p.text = f"{s_pct:.0f}%"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _score_color_rgb(s_pct)
    p.alignment = PP_ALIGN.RIGHT

    # Stats line (computed across the whole audit, shown on every slide)
    passed, failed, skipped, total = stats
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.3))
    p = stats_box.text_frame.paragraphs[0]
    p.text = (
        f"{passed} PASS  •  {failed} FAIL  •  {skipped} SKIP  •  {total} Total"
        f"   |   Score: {s_passed}/{s_total}"
    )
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 100, 100)

    # Table: Status | Check | Severity | Detail
    rows = len(checks) + 1
    table = slide.shapes.add_table(
        rows, 4, Inches(0.4), Inches(1.4), Inches(9.2), Inches(0.4 * rows)
    ).table
    table.columns[0].width = Inches(0.9)   # status
    table.columns[1].width = Inches(2.6)   # check name
    table.columns[2].width = Inches(1.1)   # severity
    table.columns[3].width = Inches(4.6)   # detail

    # Header row
    for col, label in enumerate(("Status", "Check", "Severity", "Detail")):
        _set_cell(
            table.cell(0, col), label,
            size=11, bold=True, color=_HEADER_TEXT, fill=_HEADER_FILL,
        )

    # Data rows
    for i, check in enumerate(checks, start=1):
        status_text, status_color = _status_of(check)
        row_fill = _ROW_FILL_A if i % 2 else _ROW_FILL_B
        sev = check.severity.value
        detail = check.detail.replace("\n", " ").strip()
        if len(detail) > _MAX_DETAIL_CHARS:
            detail = detail[:_MAX_DETAIL_CHARS - 1] + "…"

        _set_cell(table.cell(i, 0), status_text, size=10, bold=True,
                  color=status_color, align=PP_ALIGN.CENTER, fill=row_fill)
        _set_cell(table.cell(i, 1), check.name, size=10, bold=True,
                  color=RGBColor(33, 33, 33), fill=row_fill)
        _set_cell(table.cell(i, 2), sev, size=9,
                  color=_SEVERITY_COLOR.get(sev, RGBColor(100, 100, 100)),
                  align=PP_ALIGN.CENTER, fill=row_fill)
        _set_cell(table.cell(i, 3), detail, size=9, color=_DETAIL_TEXT, fill=row_fill)


def _add_audit_slides(
    prs: Presentation, journey_name: str, checks: List[CheckResult]
) -> None:
    """Add one or more slides for an audit, paginating so no check is dropped."""
    passed = sum(1 for c in checks if c.passed and not c.skipped)
    failed = sum(1 for c in checks if not c.passed and not c.skipped)
    skipped = sum(1 for c in checks if c.skipped)
    stats = (passed, failed, skipped, len(checks))
    score = score_checks(checks)

    if not checks:
        _add_audit_slide(prs, journey_name, [], stats, score)
        return

    chunks = [
        checks[i:i + _ROWS_PER_SLIDE]
        for i in range(0, len(checks), _ROWS_PER_SLIDE)
    ]
    for idx, chunk in enumerate(chunks):
        title = journey_name if idx == 0 else f"{journey_name}  (cont.)"
        _add_audit_slide(prs, title, chunk, stats, score)


def _add_summary_slide(
    prs: Presentation,
    journey_results: List[Tuple[str, List[CheckResult]]],
) -> None:
    """Scorecard slide: per-audit equal-weighted score plus an overall score."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Audit Scorecard"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(33, 33, 33)

    # Table: Audit | Score | Checks (one row per audit + header + overall)
    rows = len(journey_results) + 2
    table = slide.shapes.add_table(
        rows, 3, Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.45 * rows)
    ).table
    table.columns[0].width = Inches(5.0)   # audit name
    table.columns[1].width = Inches(1.9)   # score %
    table.columns[2].width = Inches(1.9)   # passed/total

    for col, label in enumerate(("Audit", "Score", "Checks")):
        _set_cell(table.cell(0, col), label, size=12, bold=True,
                  color=_HEADER_TEXT, fill=_HEADER_FILL)

    overall_passed = overall_total = 0
    for i, (journey_name, checks) in enumerate(journey_results, start=1):
        passed, total, pct = score_checks(checks)
        overall_passed += passed
        overall_total += total
        row_fill = _ROW_FILL_A if i % 2 else _ROW_FILL_B
        _set_cell(table.cell(i, 0), journey_name, size=12, bold=True,
                  color=RGBColor(33, 33, 33), fill=row_fill)
        _set_cell(table.cell(i, 1), f"{pct:.0f}%", size=12, bold=True,
                  color=_score_color_rgb(pct), align=PP_ALIGN.CENTER, fill=row_fill)
        _set_cell(table.cell(i, 2), f"{passed}/{total}", size=11,
                  color=_DETAIL_TEXT, align=PP_ALIGN.CENTER, fill=row_fill)

    # Overall row
    overall_pct = (overall_passed / overall_total * 100.0) if overall_total else 0.0
    last = rows - 1
    _set_cell(table.cell(last, 0), "OVERALL", size=12, bold=True,
              color=_HEADER_TEXT, fill=_HEADER_FILL)
    _set_cell(table.cell(last, 1), f"{overall_pct:.0f}%", size=12, bold=True,
              color=_HEADER_TEXT, align=PP_ALIGN.CENTER, fill=_score_color_rgb(overall_pct))
    _set_cell(table.cell(last, 2), f"{overall_passed}/{overall_total}", size=11, bold=True,
              color=_HEADER_TEXT, align=PP_ALIGN.CENTER, fill=_HEADER_FILL)


def export_to_powerpoint(
    journey_results: List[Tuple[str, List[CheckResult]]],
    site_url: str,
    output_path: str,
) -> None:
    """
    Generate a PowerPoint presentation from audit results.

    Args:
        journey_results: list of (journey_name, checks) tuples
        site_url: the site URL being audited
        output_path: where to save the .pptx file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide, then scorecard summary
    _add_title_slide(prs, site_url)
    _add_summary_slide(prs, journey_results)

    # One or more slides per audit (paginated so nothing is dropped)
    for journey_name, checks in journey_results:
        _add_audit_slides(prs, journey_name, checks)

    # Save
    prs.save(output_path)
    print(f"✓ PowerPoint report saved to: {output_path}")
